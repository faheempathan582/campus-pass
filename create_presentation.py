import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    BG_DARK = RGBColor(6, 11, 24)        # #060B18 Deep Navy
    CARD_BG = RGBColor(17, 24, 39)       # #111827 Dark Surface
    CARD_BORDER = RGBColor(45, 55, 72)   # #2D3748 Border
    PRIMARY = RGBColor(108, 99, 255)     # #6C63FF Electric Indigo
    PRIMARY_LIGHT = RGBColor(165, 180, 252) # #A5B4FC
    TEXT_MAIN = RGBColor(240, 244, 255)  # #F0F4FF White/Blue
    TEXT_MUTED = RGBColor(148, 163, 184) # #94A3B8 Muted Gray
    ACCENT_GREEN = RGBColor(16, 185, 129)# #10B981 Success
    ACCENT_AMBER = RGBColor(245, 158, 11) # #F59E0B Warning
    ACCENT_PINK = RGBColor(244, 114, 182) # #F472B6 Pink Accent

    blank_slide_layout = prs.slide_layouts[6]

    def set_slide_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="CAMPUS PASS TECHNICAL DECK"):
        # Category tag
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_LIGHT
        p.font.name = "Calibri"

        # Title text
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(26)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_MAIN
        p2.font.name = "Calibri"

    def add_card(slide, left, top, width, height, title, content_list, accent_color=PRIMARY):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1)

        # Top accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_color
        bar.line.fill.background()

        # Text Frame
        txBox = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(width - 0.4), Inches(height - 0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.1)
        tf.margin_left = Inches(0.1)

        if title:
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = TEXT_MAIN
            p.font.name = "Calibri"
            p.space_after = Pt(10)
            first_bullet = True
        else:
            first_bullet = False

        for idx, item in enumerate(content_list):
            if first_bullet and idx == 0:
                p = tf.add_paragraph()
            elif not title and idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            if isinstance(item, tuple):
                header, desc = item
                p.text = f"•  {header}: "
                p.font.bold = True
                p.font.size = Pt(12)
                p.font.color.rgb = PRIMARY_LIGHT
                p.font.name = "Calibri"
                
                # Add normal text inline if possible, or string format
                run = p.add_run()
                run.text = desc
                run.font.bold = False
                run.font.color.rgb = TEXT_MUTED
            else:
                p.text = f"•  {item}"
                p.font.size = Pt(12)
                p.font.color.rgb = TEXT_MUTED
                p.font.name = "Calibri"
            p.space_after = Pt(6)

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s1)

    # Title Banner Box
    tb = s1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "CAMPUSPASS"
    p0.font.size = Pt(48)
    p0.font.bold = True
    p0.font.color.rgb = PRIMARY
    p0.font.name = "Calibri"

    p1 = tf.add_paragraph()
    p1.text = "Next-Gen Digital Permission & Security Platform"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_MAIN
    p1.font.name = "Calibri"
    p1.space_after = Pt(20)

    p2 = tf.add_paragraph()
    p2.text = "Technical Architecture, Security Hardening & Implementation Overview"
    p2.font.size = Pt(16)
    p2.font.color.rgb = TEXT_MUTED
    p2.font.name = "Calibri"
    p2.space_after = Pt(40)

    p3 = tf.add_paragraph()
    p3.text = "Target Institution: SRIT (@srit.ac.in)   |   Live Demo: campus-pass-cspy.vercel.app"
    p3.font.size = Pt(13)
    p3.font.color.rgb = ACCENT_GREEN
    p3.font.name = "Calibri"

    # -------------------------------------------------------------
    # SLIDE 2: Executive Summary & System Objectives
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s2)
    add_header(s2, "Executive Summary & Core System Objectives")

    add_card(s2, 0.8, 1.6, 5.6, 5.2, "🎯 Problem & Solution", [
        ("Legacy Issue", "Manual paper leave slips caused delays, loss of records, and unauthorized access from non-college emails."),
        ("Modern Solution", "CampusPass digitizes student leave, on-duty, and hostel exit applications with real-time tracking."),
        ("Domain Lockout", "Strictly restricts registration and authentication to verified @srit.ac.in emails."),
        ("Multi-Tier Approvals", "Automated routing through Class Advisor -> HOD -> Hostel Warden workflows.")
    ], PRIMARY)

    add_card(s2, 6.8, 1.6, 5.6, 5.2, "🛡️ Key Security Deliverables", [
        ("Domain Validation", "Live regex checking on Frontend + strict validation on Backend endpoints."),
        ("Protected Routes", "JWT-based client-side route guards blocking unauthorized URL navigation."),
        ("Logout Management", "Dynamic auth state dispatch with immediate token revocation."),
        ("CORS Hardening", "Restricted API access strictly to the Vercel frontend production domain.")
    ], ACCENT_GREEN)

    # -------------------------------------------------------------
    # SLIDE 3: Complete Technology Stack
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s3)
    add_header(s3, "Full System Technology Stack")

    add_card(s3, 0.8, 1.6, 3.6, 5.2, "⚡ Frontend Layer", [
        ("Framework", "React 19 + Vite 8"),
        ("Routing", "React Router DOM v7"),
        ("Icons", "Lucide React Vector Library"),
        ("Styling", "Custom Glassmorphic CSS"),
        ("Typography", "Google Font Outfit")
    ], PRIMARY)

    add_card(s3, 4.8, 1.6, 3.6, 5.2, "⚙️ Backend & API", [
        ("Runtime", "Node.js (v20+)"),
        ("Framework", "Express.js REST Server"),
        ("Auth Engine", "JWT (JSON Web Tokens)"),
        ("Hashing", "Bcrypt.js (10 salt rounds)"),
        ("Security", "CORS + Custom Middleware")
    ], ACCENT_AMBER)

    add_card(s3, 8.8, 1.6, 3.7, 5.2, "🗄️ Database & Cloud", [
        ("Primary DB", "MongoDB Cloud / Atlas"),
        ("ODM", "Mongoose Document Mapping"),
        ("Dev Fallback", "MongoMemoryServer In-Memory"),
        ("Hosting Platform", "Vercel Cloud Network"),
        ("Repository", "GitHub Continuous Deployment")
    ], ACCENT_PINK)

    # -------------------------------------------------------------
    # SLIDE 4: Security Architecture & Domain Restriction
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s4)
    add_header(s4, "Security Architecture & Domain Restrictions")

    add_card(s4, 0.8, 1.6, 5.6, 5.2, "🔒 Strict Domain Validation (@srit.ac.in)", [
        ("Frontend Guard", "Real-time input validation in Login.jsx & Register.jsx disables button & shows red alert if email does not end with @srit.ac.in."),
        ("Backend Defense", "Server-side check in validateEmail.js rejects unauthorized payloads with HTTP 400 even if API is called directly."),
        ("Security Impact", "Prevents external spam, unauthorized registrations, and potential impersonation attacks.")
    ], ACCENT_GREEN)

    add_card(s4, 6.8, 1.6, 5.6, 5.2, "🛡️ Route Protection & Session Security", [
        ("ProtectedRoute Component", "Interceptors check stored JWT and verify user roles (Student vs Advisor/HOD/Warden) before rendering dashboard components."),
        ("GuestRoute Component", "Prevents authenticated users from seeing Login/Register forms, auto-redirecting to active dashboards."),
        ("Header Authentication", "All protected API calls inject 'Authorization: Bearer <TOKEN>' header. Middleware returns HTTP 401 on invalid/expired sessions.")
    ], PRIMARY)

    # -------------------------------------------------------------
    # SLIDE 5: Database Schema & Data Models
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s5)
    add_header(s5, "Database Schemas & Data Storage Structure")

    add_card(s5, 0.8, 1.6, 3.6, 5.2, "👤 User Collection", [
        ("name", "String (Required)"),
        ("email", "Unique String (@srit.ac.in)"),
        ("password", "Bcrypt Hash String"),
        ("role", "Student / Advisor / HOD / Warden"),
        ("rollNumber", "Student Roll ID"),
        ("department", "Engineering Dept"),
        ("isHosteller", "Boolean Flag")
    ], PRIMARY)

    add_card(s5, 4.8, 1.6, 3.6, 5.2, "📄 Permission Collection", [
        ("student", "Ref to User ObjectID"),
        ("type", "Leave / On-Duty / Exit / Medical"),
        ("reason", "Text description"),
        ("fromDate / toDate", "Date Timestamps"),
        ("status", "Pending / Approved / Rejected"),
        ("pendingWithRole", "Current Approver Role"),
        ("remarks", "Array of Audit Trail Records")
    ], ACCENT_AMBER)

    add_card(s5, 8.8, 1.6, 3.7, 5.2, "🔔 Notification Collection", [
        ("recipientRole", "Target Authority Role"),
        ("message", "System Alert Description"),
        ("permission", "Ref to Permission ObjectID"),
        ("readBy", "Array of User ObjectIDs"),
        ("timestamps", "CreatedAt & UpdatedAt")
    ], ACCENT_GREEN)

    # -------------------------------------------------------------
    # SLIDE 6: Approval Workflow & Notification Engine
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s6)
    add_header(s6, "Approval Workflow & Real-Time Notification Engine")

    add_card(s6, 0.8, 1.6, 11.7, 5.2, "🔄 Multi-Level Approval Sequence", [
        ("Step 1 - Application Submission", "Student completes permission form. Backend generates permission document setting pendingWithRole to 'Advisor' (or 'Warden' for Hostel Exit)."),
        ("Step 2 - Initial Authority Notification", "System auto-generates a Notification document tagged for the recipient role. NotificationBell polls every 30 seconds to alert faculty."),
        ("Step 3 - First Level Review", "Class Advisor reviews request in AuthorityDashboard, adds optional remark, and clicks Approve or Reject."),
        ("Step 4 - Escalation & Final Approval", "If approved by Advisor, status stays 'Pending' while pendingWithRole advances to 'HOD' & triggers HOD notification. When HOD approves, status updates to 'Approved'."),
        ("Step 5 - Student History Sync", "Student Dashboard updates live with real-time status badges (Pending, Approved, Rejected) and remarks.")
    ], ACCENT_PINK)

    # -------------------------------------------------------------
    # SLIDE 7: UI Redesign & Design System
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s7)
    add_header(s7, "Frontend Redesign & User Experience (UX)")

    add_card(s7, 0.8, 1.6, 5.6, 5.2, "🎨 Modern Design System", [
        ("Deep Navy Palette", "Rich dark mode aesthetic (#060b18) featuring subtle radial mesh backgrounds."),
        ("Glassmorphism", "Translucent cards with blur filters, subtle inset shadows, and reactive hover borders."),
        ("Typography & Icons", "Google Font 'Outfit' combined with clean vector icons from Lucide React.")
    ], PRIMARY)

    add_card(s7, 6.8, 1.6, 5.6, 5.2, "💡 Enhanced UI Components", [
        ("Custom Toast System", "Toast.jsx replaces native browser alert() popups with smooth non-blocking notifications."),
        ("Interactive Dashboards", "Student stats cards (Total, Pending, Approved) + expandable Authority review cards."),
        ("Responsive Navigation", "Sticky frosted navbar with live user avatar chip and red accent Logout button.")
    ], ACCENT_AMBER)

    # -------------------------------------------------------------
    # SLIDE 8: Deployment, Build Verification & Summary
    # -------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s8)
    add_header(s8, "Production Build & Live Deployment Summary")

    add_card(s8, 0.8, 1.6, 5.6, 5.2, "🚀 Continuous Deployment Pipeline", [
        ("Local Vite Build", "Tested via 'npx vite build'. Successfully compiled 1,794 modules with zero errors in 1.62 seconds."),
        ("Git Version Control", "Pushed 20 modified/created files directly to GitHub main branch."),
        ("Vercel Auto-Deploy", "Triggered instant CDN build serving live application at campus-pass-cspy.vercel.app.")
    ], ACCENT_GREEN)

    add_card(s8, 6.8, 1.6, 5.6, 5.2, "✅ Project Summary", [
        ("Email Restrict", "Strictly enforced @srit.ac.in on Frontend + Backend."),
        ("Authentication", "Navbar Logout + JWT Protected Route Guards."),
        ("Notifications", "Automated alert engine for Advisor, HOD, and Warden."),
        ("Live Status", "Fully operational and available for institution testing.")
    ], PRIMARY)

    # Save presentation
    output_path = r"C:\Users\DELL\.gemini\antigravity\scratch\CampusPass\CampusPass.pptx"
    prs.save(output_path)
    print(f"SUCCESS: Created presentation at {output_path}")

if __name__ == "__main__":
    create_presentation()
